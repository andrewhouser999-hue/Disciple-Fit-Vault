#!/usr/bin/env python3
"""
Build A001 — Presentation Deck 3 (1-1 Conversational Backup)
6 slides, white bg, navy/teal, no bullet points, speaker notes.
Output: A001 — Presentation Deck 3.pptx  (same folder)
Run: python build_deck3.py
"""

import os
import zipfile
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

OUTPUT = r"C:\Users\Andre\Documents\Obsidian Vault\Andy's Life\Coaching\Business Advisory Services\A001\A001 — Presentation Deck 3.pptx"

NAVY  = RGBColor(0x1A, 0x27, 0x44)
TEAL  = RGBColor(0x2E, 0x7D, 0x87)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0xCC, 0xCC, 0xCC)
FONT  = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK  = prs.slide_layouts[6]
MAR    = Inches(0.5)
USE_W  = prs.slide_width - 2 * MAR   # 12.33"
W      = prs.slide_width
H      = prs.slide_height


# ── Helpers ───────────────────────────────────────────────────────────────────

def white_bg(slide):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = WHITE


def tb(slide, l, t, w, h, text,
       size=24, bold=False, italic=False, color=NAVY,
       align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    p   = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text        = text
    run.font.name   = FONT
    run.font.size   = Pt(size)
    run.font.bold   = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box


def tb_ml(slide, l, t, w, h, lines,
          size=22, bold=False, italic=False, color=NAVY,
          align=PP_ALIGN.CENTER):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf  = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text        = line
        run.font.name   = FONT
        run.font.size   = Pt(size)
        run.font.bold   = bold
        run.font.italic = italic
        run.font.color.rgb = color
    return box


def rect(slide, l, t, w, h, fill=WHITE, border=TEAL, bpt=1.5):
    s = slide.shapes.add_shape(1, l, t, w, h)  # 1 = rectangle
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = border
    s.line.width = Pt(bpt)
    return s


def box_text(slide, l, t, w, h, heading, lines,
             hsize=20, bsize=18):
    s  = rect(slide, l, t, w, h)
    tf = s.text_frame
    tf.word_wrap    = True
    tf.margin_left  = Inches(0.08)
    tf.margin_right = Inches(0.08)
    tf.margin_top   = Inches(0.12)

    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = heading
    run.font.name  = FONT
    run.font.size  = Pt(hsize)
    run.font.bold  = True
    run.font.color.rgb = TEAL

    for line in lines:
        p   = tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = line
        run.font.name  = FONT
        run.font.size  = Pt(bsize)
        run.font.color.rgb = NAVY
    return s


def arrow_r(slide, l, t, w, h, color=TEAL):
    s = slide.shapes.add_shape(13, l, t, w, h)  # 13 = right arrow
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.color.rgb = color
    return s


def hline(slide, l, t, w, color=LGRAY):
    s = slide.shapes.add_shape(1, l, t, w, Inches(0.025))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def vline(slide, x, t, h, color=LGRAY):
    s = slide.shapes.add_shape(1, x, t, Inches(0.025), h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — The Gap
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
white_bg(s1)

COL_W   = Inches(5.5)
GAP_W   = Inches(1.33)
C1      = MAR
C2      = MAR + COL_W + GAP_W
CTOP    = Inches(0.9)

# vertical divider first (behind arrow)
vline(s1, C1 + COL_W + GAP_W / 2, CTOP, Inches(5.0))

# arrow over divider
arrow_r(s1,
        C1 + COL_W + Inches(0.17),
        Inches(3.3),
        GAP_W - Inches(0.34),
        Inches(0.45))

# Column 1
tb(s1, C1, CTOP, COL_W, Inches(0.55),
   "STANDARD APPROACH", size=26, bold=True, align=PP_ALIGN.CENTER)
tb_ml(s1, C1, CTOP + Inches(0.7), COL_W, Inches(4.2),
      ["Changes knowledge",
       "↓",
       "Doesn't reach what produces the behavior",
       "↓",
       "Reverts in 30–90 days"],
      size=22)

# Column 2
tb(s1, C2, CTOP, COL_W, Inches(0.55),
   "WHAT'S DIFFERENT HERE", size=26, bold=True, align=PP_ALIGN.CENTER)
tb_ml(s1, C2, CTOP + Inches(0.7), COL_W, Inches(4.2),
      ["Changes conditions",
       "↓",
       "Addresses what produces the behavior",
       "↓",
       "Holds"],
      size=22)

# Bottom caption
tb(s1, MAR, Inches(6.65), USE_W, Inches(0.6),
   "The pattern is produced by the environment. Standard training doesn't touch the environment.",
   size=18, italic=True, align=PP_ALIGN.CENTER)

notes(s1,
    '"You\'ve already named this — you identified at [coffee shop] that the behavior is coming from the '
    'environment, not the character of your people. That\'s exactly right, and it\'s also why what '
    'you\'ve tried hasn\'t held. Standard training changes what people know. It doesn\'t touch the '
    'conditions that produced the behavior in the first place. When your manager walks back into the '
    'same environment, the same conditions produce the same result. What\'s different here is the '
    'level of the intervention. We work where the pattern is actually made."\n\n'
    'Deploy when: Owner mentions prior training that failed, OR asks "why does this keep happening?"\n'
    'Duration: 60–75 seconds\n'
    'Put away when: Owner acknowledges the distinction — any verbal signal that it landed')


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Chain
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
white_bg(s2)

tb(s2, MAR, Inches(0.4), USE_W, Inches(0.65),
   "Why the same patterns keep showing up",
   size=34, bold=True, align=PP_ALIGN.LEFT)

ARROW_W  = Inches(0.28)
BOX_TOP  = Inches(1.3)
BOX_H    = Inches(2.1)
N        = 5
box_w    = (USE_W - 4 * ARROW_W) / N

BOX_LABELS = [
    "Structural\nConditions",
    "Manager\nExperience",
    "Employee\nExperience",
    "What\nYou See",
    "What\nIt Costs",
]
SUBTEXTS = {
    3: "conflict · avoidance\nturnover · inconsistency",
    4: "$840K–$1.6M/year",
}

cur = MAR
for i, label in enumerate(BOX_LABELS):
    b  = rect(s2, cur, BOX_TOP, box_w, BOX_H, bpt=2.0)
    tf = b.text_frame
    tf.word_wrap   = True
    tf.margin_top  = Inches(0.4)
    tf.margin_left = Inches(0.05)
    tf.margin_right= Inches(0.05)
    p   = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.name  = FONT
    run.font.size  = Pt(19)
    run.font.bold  = True
    run.font.color.rgb = NAVY

    if i in SUBTEXTS:
        tb(s2, cur, BOX_TOP + BOX_H + Inches(0.1), box_w, Inches(0.65),
           SUBTEXTS[i], size=15, italic=True, align=PP_ALIGN.CENTER)

    cur += box_w
    if i < N - 1:
        arrow_r(s2, cur, BOX_TOP + BOX_H / 2 - Inches(0.18),
                ARROW_W, Inches(0.36))
        cur += ARROW_W

tb(s2, MAR, Inches(6.65), USE_W, Inches(0.6),
   "Change the first three boxes. The last two change with them.",
   size=18, italic=True, align=PP_ALIGN.CENTER)

notes(s2,
    '"This is why the same situations keep producing the same results. Structural conditions in the '
    'operation — how authority flows, how mistakes get handled, how much clarity people have — shape '
    'how your managers experience the environment every day. That shapes how your managers show up '
    'with their teams. That shapes what your employees experience. And what you\'re watching is in '
    'the fourth box. What it\'s costing you is the fifth. You\'ve been focused on boxes four and '
    'five. The work is in the first three."\n\n'
    'Deploy when: Owner asks HOW environment produces behavior; OR after describing a specific incident\n'
    'Duration: 75–90 seconds\n'
    'Put away when: Owner can place what he described in the chain')


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — What It Costs
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
white_bg(s3)

tb(s3, MAR, Inches(0.65), USE_W, Inches(0.65),
   "400 employees   ×   ~75% annual turnover   =   ~300 replacements / year",
   size=26, align=PP_ALIGN.CENTER)

tb(s3, MAR, Inches(1.55), USE_W, Inches(1.4),
   "$1,200,000",
   size=64, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

tb(s3, MAR, Inches(3.05), USE_W, Inches(0.65),
   "in annual replacement cost — before productivity loss, training drag, and guest experience",
   size=20, align=PP_ALIGN.CENTER)

hline(s3, MAR, Inches(3.95), USE_W)

tb(s3, MAR, Inches(4.15), USE_W, Inches(0.6),
   "Phase 1 investment: $28,000",
   size=28, bold=True, align=PP_ALIGN.CENTER)

tb(s3, MAR, Inches(4.85), USE_W, Inches(0.6),
   "Recovered at a fraction of a percentage point of turnover reduction.",
   size=22, align=PP_ALIGN.CENTER)

notes(s3,
    '"At your scale, roughly 300 replacements a year. At $3,000 to $5,000 per replacement — that\'s '
    'the conservative range for hospitality — you\'re at $1.2 million annually, before you factor in '
    'the productivity gap while a position is vacant, the ramp time on a new hire, or what constant '
    'turnover does to the guest experience. The Phase 1 investment is $28,000. If turnover drops even '
    'two or three percent, it\'s paid for before you\'d feel it. The question isn\'t whether you can '
    'afford to do this. You\'ve already been paying for not doing it."\n\n'
    'Deploy when: Owner raises cost or asks about pricing; OR making the investment case before Phase 1\n'
    'Duration: 60 seconds\n'
    'Put away when: Owner has reacted to the numbers')


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — What Phase 1 Produces
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
white_bg(s4)

tb(s4, MAR, Inches(0.4), USE_W, Inches(0.65),
   "At the end of Phase 1, you have three things.",
   size=34, bold=True, align=PP_ALIGN.LEFT)

GAP3  = Inches(0.2)
BW3   = (USE_W - 2 * GAP3) / 3

BOXES = [
    ("BEHAVIORAL PROFILE",
     ["Per leader assessed",
      "What their patterns are",
      "What activates them",
      "Their development path"]),
    ("TEAM MAP",
     ["How the management team",
      "functions collectively",
      "What structural conditions",
      "are contributing"]),
    ("STRUCTURAL\nRECOMMENDATIONS",
     ["Specific changes to how",
      "the environment is designed",
      "Connected to what the",
      "assessment found"]),
]

cur = MAR
for heading, lines in BOXES:
    box_text(s4, cur, Inches(1.3), BW3, Inches(4.8),
             heading, lines, hsize=20, bsize=18)
    cur += BW3 + GAP3

tb(s4, MAR, Inches(6.65), USE_W, Inches(0.6),
   "You see the picture. You decide whether to continue.",
   size=18, italic=True, align=PP_ALIGN.CENTER)

notes(s4,
    '"Three deliverables at the end of Phase 1. A behavioral profile for each leader — what their '
    'patterns are, what activates them, what development would actually move them. A team map — how '
    'this management group functions collectively, and what structural conditions are contributing to '
    'what you\'re seeing. And specific structural recommendations — not general improvements, but '
    'changes connected to what the assessment actually found. Phase 1 gives you the picture. At the '
    'end of it, you decide whether you want to move to implementation. That\'s a deliberate structure '
    '— you\'re not committing to a full engagement on trust. You\'re committing to clarity first."\n\n'
    'Deploy when: Owner asks "what do I actually get?" or "what are the deliverables?"\n'
    'Duration: 75 seconds\n'
    'Put away when: Owner has processed all three boxes — "clarity first" signals the frame landed')


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Timeline
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
white_bg(s5)

tb(s5, MAR, Inches(0.4), USE_W, Inches(0.65),
   "When you see what.",
   size=34, bold=True, align=PP_ALIGN.LEFT)

BAR_L = MAR + Inches(1.5)
BAR_W = USE_W - Inches(3.0)   # 9.33"
BAR_T = Inches(3.8)
BAR_H = Inches(0.18)

b = rect(s5, BAR_L, BAR_T, BAR_W, BAR_H, fill=TEAL, border=TEAL, bpt=0)
b.line.fill.background()

POINTS = [
    ("Week 13",     "PHASE 1 COMPLETE",
     "Assessment findings +\nstructural recommendations delivered"),
    ("60–90 Days",  "FIRST BEHAVIORAL SIGNALS",
     "Managers functioning differently\nunder pressure\n(if implementation proceeds)"),
    ("6–12 Months", "TURNOVER SHIFT",
     "Measurable reduction\nin departure rate"),
    ("18+ Months",  "CULTURE THAT HOLDS",
     "Norms self-transmitting to new\nemployees without you\nmaintaining it manually"),
]

LW = Inches(2.5)

for i, (tlabel, title, desc) in enumerate(POINTS):
    px = BAR_L + i * BAR_W / (len(POINTS) - 1)
    ll = max(MAR, min(W - MAR - LW, px - LW / 2))

    # tick mark
    tick = s5.shapes.add_shape(1,
        px - Inches(0.035), BAR_T - Inches(0.22),
        Inches(0.07), BAR_H + Inches(0.44))
    tick.fill.solid()
    tick.fill.fore_color.rgb = NAVY
    tick.line.fill.background()

    if i % 2 == 0:          # above bar
        yt, yh, yd = Inches(1.15), Inches(1.7), Inches(2.2)
    else:                   # below bar
        yt, yh, yd = Inches(4.25), Inches(4.75), Inches(5.25)

    tb(s5, ll, yt, LW, Inches(0.4),
       tlabel, size=18, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
    tb(s5, ll, yh, LW, Inches(0.42),
       title, size=15, bold=True, align=PP_ALIGN.CENTER)
    tb(s5, ll, yd, LW, Inches(0.85),
       desc, size=14, align=PP_ALIGN.CENTER)

tb(s5, MAR, Inches(6.65), USE_W, Inches(0.6),
   "This isn't fast. What's faster is not doing it — and you've already seen what that costs.",
   size=18, italic=True, align=PP_ALIGN.CENTER)

notes(s5,
    '"Phase 1 is 13 weeks — assessment, interpretation, structural recommendations delivered. If you '
    'proceed to implementation, the first behavioral signals typically show up within 60 to 90 days '
    '— your managers functioning differently under pressure. Turnover shift follows at 6 to 12 months. '
    'A culture that sustains itself without you manually maintaining it — 18 months or more. I\'m not '
    'going to tell you this is fast. What I will say is that you\'ve already seen what fast gets you. '
    'Every training program that lasted 30 days was fast. This is different in kind, not just in content."\n\n'
    'Deploy when: Owner asks "how long before I see a difference?" or expresses timeline pressure\n'
    'Duration: 75 seconds\n'
    'Put away when: Owner acknowledges the timeline — "so Phase 1 is essentially this year..." signals it landed')


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — What Your Managers Will Experience
# ══════════════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(BLANK)
white_bg(s6)

tb(s6, MAR, Inches(0.4), USE_W, Inches(0.65),
   "Before you introduce this to your team.",
   size=34, bold=True, align=PP_ALIGN.LEFT)

C1_6 = MAR
C2_6 = MAR + COL_W + GAP_W
CTOP6 = Inches(1.25)

# divider
vline(s6, C1_6 + COL_W + GAP_W / 2, CTOP6, Inches(5.0))

# Column 1
tb(s6, C1_6, CTOP6, COL_W, Inches(0.55),
   "What this is NOT", size=24, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
tb_ml(s6, C1_6, CTOP6 + Inches(0.65), COL_W, Inches(4.2),
      ["An evaluation of how they manage",
       "An HR audit or compliance review",
       "Individual data going to you",
       "Supervision by an outside person"],
      size=22)

# Column 2
tb(s6, C2_6, CTOP6, COL_W, Inches(0.55),
   "What this IS", size=24, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
tb_ml(s6, C2_6, CTOP6 + Inches(0.65), COL_W, Inches(4.2),
      ["Private development — their data stays with them",
       "A picture of how the team functions collectively",
       "Tools to do their jobs better",
       "Built with them, not done to them"],
      size=22)

tb(s6, MAR, Inches(6.65), USE_W, Inches(0.6),
   "Your managers are the most important part of what we build.",
   size=18, italic=True, align=PP_ALIGN.CENTER)

notes(s6,
    '"When you bring this to your OMs, the first thing they\'ll want to know is what it means for '
    'them — whether they\'re being assessed, whether you\'ll have access to what they share. Here\'s '
    'the answer: their individual data stays with them. What you receive is the aggregate picture — '
    'how the team functions collectively and what the structure is contributing. That confidentiality '
    'isn\'t a courtesy — it\'s what makes honest assessment possible. An OM who\'s guarding what he '
    'says won\'t give you usable data. The protection is structural, not just a promise. Your managers '
    'are the lever here. They\'re not the problem — and how you introduce this will either confirm '
    'that or undo it."\n\n'
    'Deploy when: Close of meeting — owner surfaces OM question; OR Pivot 2 is raised\n'
    'Duration: 75–90 seconds\n'
    'Put away when: Owner has language for how to introduce Andrew to OMs — "so I\'d say something like..." signals the frame landed')


# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
prs.save(OUTPUT)
print(f"Saved: {OUTPUT}")


def expand_notes_panel(filepath, slide_pct=62):
    """Patch presProps.xml so the notes panel opens expanded in Normal view.
    slide_pct: percentage of the view height given to the slide area (rest = notes).
    """
    P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
    sz = str(slide_pct * 1000)   # OOXML unit: thousandths of a percent
    tmp = filepath + ".tmp"

    with zipfile.ZipFile(filepath, "r") as zin, \
         zipfile.ZipFile(tmp, "w", zipfile.ZIP_DEFLATED) as zout:
        for item in zin.infolist():
            data = zin.read(item.filename)
            if item.filename == "ppt/presProps.xml":
                root = etree.fromstring(data)
                nv = root.find(f"{{{P_NS}}}normalViewPr")
                if nv is None:
                    nv = etree.SubElement(root, f"{{{P_NS}}}normalViewPr")
                # restoredLeft must come before restoredTop per schema
                rl = nv.find(f"{{{P_NS}}}restoredLeft")
                if rl is None:
                    rl = etree.SubElement(nv, f"{{{P_NS}}}restoredLeft")
                rl.set("sz", "15620")
                rl.set("autoAdjust", "0")
                rt = nv.find(f"{{{P_NS}}}restoredTop")
                if rt is None:
                    rt = etree.SubElement(nv, f"{{{P_NS}}}restoredTop")
                rt.set("sz", sz)
                rt.set("autoAdjust", "0")
                data = etree.tostring(root, xml_declaration=True,
                                      encoding="UTF-8", standalone=True)
            zout.writestr(item, data)

    os.replace(tmp, filepath)
    print("Notes panel expanded.")


expand_notes_panel(OUTPUT)
